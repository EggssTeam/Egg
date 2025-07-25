import base64
import io
import os
from typing import List

import fitz
import requests
import re  # 用于正则表达式解析文本

from fastapi import UploadFile
from pdfminer.high_level import extract_text as extract_pdf_text
from pptx import Presentation
import tempfile

from config import DEEPSEEK_API_KEY, DEEPSEEK_API_URL, GPT4O_API_URL, OPENAI_API_KEY

if not DEEPSEEK_API_KEY:
    raise RuntimeError("环境变量 DEEPSEEK_API_KEY 未设置，请先配置！")



def extract_images_from_pdf(file_stream: io.BytesIO) -> list[bytes]:
    images = []
    doc = fitz.open(stream=file_stream, filetype="pdf")
    for page_index in range(len(doc)):
        page = doc[page_index]
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            images.append(image_bytes)
    return images


def extract_images_from_pptx(file_stream: io.BytesIO) -> list[bytes]:
    prs = Presentation(file_stream)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture 类型
                image = shape.image
                images.append(image.blob)
    return images
#

def analyze_images_with_o4(images: List[bytes]) -> str:
    descriptions = []
    for i, image in enumerate(images):
        messages = [
            {
                "role": "system",
                "content": "你是一个图像内容分析助手，请简洁描述图像内容，输出格式为：图片X：内容"
            },
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": f"请分析图像 {i + 1} 的内容。"},
                    {"type": "image_url",
                     "image_url": {"url": f"data:image/jpeg;base64,{base64.b64encode(image).decode()}"}}
                ]
            }
        ]
        response = requests.post(
            GPT4O_API_URL,
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "gpt-4-vision-preview",
                "messages": messages,
                "max_tokens": 600
            }
        )
        response.raise_for_status()
        result = response.json()
        content = result["choices"][0]["message"]["content"].strip()
        descriptions.append(f"图片{i + 1}：{content}")

    return "\n".join(descriptions)



def extract_text_from_file(file_or_bytes, suffix=None) -> str:
    import io
    from pptx import Presentation
    from pdfminer.high_level import extract_text as extract_pdf_text
    import tempfile

    if isinstance(file_or_bytes, UploadFile):
        suffix = file_or_bytes.filename.lower().split('.')[-1]
        content = file_or_bytes.file.read()
    elif isinstance(file_or_bytes, (bytes, io.BytesIO)):
        if suffix is None:
            raise ValueError("必须提供 suffix 参数")
        if isinstance(file_or_bytes, io.BytesIO):
            content = file_or_bytes.read()
        else:
            content = file_or_bytes
    else:
        raise TypeError("参数类型错误，必须是 UploadFile、bytes 或 BytesIO")

    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{suffix}") as temp:
        temp.write(content)
        temp_path = temp.name

    if suffix == 'pdf':
        return extract_pdf_text(temp_path)
    elif suffix == 'pptx':
        prs = Presentation(temp_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError("不支持的文件格式，仅支持 PDF 和 PPTX")








def generate_questions(text: str):
    prompt = f"""
请根据以下文本生成 **5 个选择题**，每个问题包含 **4 个选项（A、B、C、D）**，并标注正确答案。格式如下：

问题1: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 仅给出字母

问题2: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 仅给出字母

问题3: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 仅给出字母

问题4: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 仅给出字母

问题5: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 仅给出字母

文本内容:
{text}
"""

    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "system", "content": "你是一个擅长出选择题的AI助手，问题清晰，选项合理。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7
    }

    try:
        response = requests.post(DEEPSEEK_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        response_data = response.json()

        if "choices" not in response_data or not response_data["choices"]:
            print("API 返回数据没有包含 'choices' 或返回为空")
            return []

        qa_text = response_data["choices"][0]["message"]["content"]
        lines = [line.strip() for line in qa_text.splitlines() if line.strip()]

        questions = []
        i = 0
        while i < len(lines):
            if lines[i].startswith("问题"):
                try:
                    question_text = lines[i].split(":", 1)[1].strip()
                    options = {}
                    for j in range(1, 5):
                        opt_line = lines[i + j]
                        if ")" not in opt_line:
                            print(f"选项格式错误: {opt_line}")
                            break
                        key = opt_line[0]
                        val = opt_line[opt_line.index(")") + 1:].strip()
                        options[key] = val
                    correct_answer_line = lines[i + 5]
                    if ":" not in correct_answer_line:
                        print(f"正确答案格式错误: {correct_answer_line}")
                        break
                    correct_answer = correct_answer_line.split(":")[1].strip()
                    questions.append({
                        "question": question_text,
                        "options": options,
                        "correct_answer": correct_answer
                    })
                    i += 6
                except Exception as e:
                    print(f"解析失败: {e}")
                    break
            else:
                i += 1

        return questions

    except Exception as e:
        print(f"API 请求失败: {e}")
        return []
