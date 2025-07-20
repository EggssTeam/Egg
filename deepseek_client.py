import os
import requests
import re  # 用于正则表达式解析文本

from fastapi import UploadFile
from pdfminer.high_level import extract_text as extract_pdf_text
from pptx import Presentation
import tempfile
from pdfminer.high_level import extract_text

DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")



if not DEEPSEEK_API_KEY:
    raise RuntimeError("环境变量 DEEPSEEK_API_KEY 未设置，请先配置！")
#
# def pdf_to_text(file_path: str) -> str:
#     return extract_text(file_path)
# #


def extract_text_from_file(file: UploadFile) -> str:
    suffix = file.filename.lower().split('.')[-1]
    content = file.file.read()  # 读取二进制内容

    # 写入临时文件以供解析
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{suffix}") as temp:
        temp.write(content)
        temp_path = temp.name

    if suffix == 'pdf':
        return extract_pdf_text(temp_path)
    elif suffix == 'pptx':
        prs = Presentation(temp_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError("不支持的文件格式，仅支持 PDF 和 PPTX")




def generate_questions(text: str):
    prompt = f"""
请根据以下文本生成 **3 个选择题**，每个问题包含 **4 个选项（A、B、C、D）**，并标注正确答案。格式如下：

问题1: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 给出字母+具体解释，不要重复选项内容

问题2: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 给出字母+具体解释，不要重复选项内容

问题3: ...
A) ...
B) ...
C) ...
D) ...
正确答案: 给出字母+具体解释，不要重复选项内容

文本内容:
{text}
"""

    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "system", "content": "你是一个擅长出选择题的AI助手，问题清晰，选项合理。"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7
    }

    try:
        response = requests.post(DEEPSEEK_API_URL, headers=headers, json=payload)
        response.raise_for_status()
        response_data = response.json()

        if "choices" not in response_data or not response_data["choices"]:
            print("API 返回数据没有包含 'choices' 或返回为空")
            return []

        qa_text = response_data["choices"][0]["message"]["content"]
        lines = [line.strip() for line in qa_text.splitlines() if line.strip()]

        questions = []
        i = 0
        while i < len(lines):
            if lines[i].startswith("问题"):
                try:
                    question_text = lines[i].split(":", 1)[1].strip()
                    options = {}
                    for j in range(1, 5):
                        opt_line = lines[i + j]
                        if ")" not in opt_line:
                            print(f"选项格式错误: {opt_line}")
                            break
                        key = opt_line[0]
                        val = opt_line[opt_line.index(")") + 1:].strip()
                        options[key] = val
                    correct_answer_line = lines[i + 5]
                    if ":" not in correct_answer_line:
                        print(f"正确答案格式错误: {correct_answer_line}")
                        break
                    correct_answer = correct_answer_line.split(":")[1].strip()
                    questions.append({
                        "question": question_text,
                        "options": options,
                        "correct_answer": correct_answer
                    })
                    i += 6
                except Exception as e:
                    print(f"解析失败: {e}")
                    break
            else:
                i += 1

        return questions

    except Exception as e:
        print(f"API 请求失败: {e}")
        return []
